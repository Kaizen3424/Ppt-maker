"""PPTX compiler - converts JSON deck to PowerPoint file."""
import os
from typing import Dict, Any, List, Optional

# Import at module level for python-pptx
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False


# Default theme configurations
DEFAULT_THEMES = {
    "default": {
        "title_font": "Arial",
        "body_font": "Arial",
        "title_size": 44,
        "heading_size": 32,
        "body_size": 18,
        "bullet_size": 16,
        "primary_color": RGBColor(102, 126, 234),  # #667eea
        "secondary_color": RGBColor(118, 75, 162),  # #764ba2
        "text_color": RGBColor(255, 255, 255),
        "background_gradient": True
    },
    "corporate": {
        "title_font": "Calibri",
        "body_font": "Calibri",
        "title_size": 40,
        "heading_size": 28,
        "body_size": 16,
        "bullet_size": 14,
        "primary_color": RGBColor(0, 51, 102),  # #003366
        "secondary_color": RGBColor(51, 102, 153),  # #336699
        "text_color": RGBColor(255, 255, 255),
        "background_gradient": False
    },
    "modern": {
        "title_font": "Segoe UI",
        "body_font": "Segoe UI",
        "title_size": 48,
        "heading_size": 36,
        "body_size": 20,
        "bullet_size": 18,
        "primary_color": RGBColor(255, 87, 51),  # #ff5733
        "secondary_color": RGBColor(255, 153, 0),  # #ff9900
        "text_color": RGBColor(255, 255, 255),
        "background_gradient": True
    },
    "minimal": {
        "title_font": "Helvetica",
        "body_font": "Helvetica",
        "title_size": 42,
        "heading_size": 30,
        "body_size": 18,
        "bullet_size": 16,
        "primary_color": RGBColor(33, 33, 33),  # #212121
        "secondary_color": RGBColor(66, 66, 66),  # #424242
        "text_color": RGBColor(0, 0, 0),
        "background_gradient": False
    }
}


def generate_pptx(json_deck: Dict[str, Any], output_path: str) -> str:
    """Generate a PowerPoint file from JSON deck data.
    
    Args:
        json_deck: The deck data with slides
        output_path: Path where the PPTX will be saved
        
    Returns:
        Path to the generated PPTX file
    """
    if not PPTX_AVAILABLE:
        # Fallback: create a placeholder
        return _create_placeholder_pptx(output_path)
    
    # Get theme from deck metadata
    theme_name = json_deck.get("theme", "default")
    theme = DEFAULT_THEMES.get(theme_name, DEFAULT_THEMES["default"])
    
    # Create presentation
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 ratio
    prs.slide_height = Inches(7.5)
    
    slides = json_deck.get("slides", [])
    
    for slide_data in slides:
        _add_slide(prs, slide_data, theme)
    
    # Save the presentation
    prs.save(output_path)
    return output_path


def _add_slide(prs: 'Presentation', slide_data: Dict[str, Any], theme: Dict[str, Any]) -> None:
    """Add a single slide to the presentation with rich content support."""
    
    slide_type = slide_data.get("type", "content")
    title = slide_data.get("title", "")
    heading = slide_data.get("heading", title)
    body_text = slide_data.get("body_text", "")
    bullets = slide_data.get("bullet_points", [])
    callout = slide_data.get("callout", "")
    
    # Get layout from slide data
    layout_spec = slide_data.get("layout", {})
    elements = layout_spec.get("elements", [])
    
    # Get background styling from layout
    background = layout_spec.get("background", {})
    
    # Get layout
    if slide_type == "title":
        layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = heading
        
        # Set subtitle if available
        if len(slide.placeholders) > 1:
            subtitle = slide.placeholders[1]
            subtitle.text = slide_data.get("subheading", body_text)
    else:
        # Content slide
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = heading
        
        # Process elements based on layout specification
        if elements:
            _add_layout_elements(slide, slide_data, theme, elements)
        else:
            # Fallback: use default content placeholder
            _add_default_content(slide, body_text, bullets, theme)
        
        # Add callout if present
        if callout:
            _add_callout_box(slide, callout, theme)
    
    # Apply background styling
    _apply_background(slide, background, theme)


def _add_layout_elements(slide, slide_data: Dict[str, Any], theme: Dict[str, Any], elements: List[Dict[str, Any]]) -> None:
    """Add elements to slide based on layout specification."""
    
    heading = slide_data.get("heading", "")
    body_text = slide_data.get("body_text", "")
    bullets = slide_data.get("bullet_points", [])
    callout = slide_data.get("callout", "")
    subheading = slide_data.get("subheading", "")
    
    # Check for tables
    tables = slide_data.get("tables", [])
    for i, table_data in enumerate(tables):
        _add_table(slide, table_data, theme)
    
    # Check for images
    images = slide_data.get("images", [])
    for image_path in images:
        _add_image(slide, image_path, theme)
    
    # Check for charts
    charts = slide_data.get("charts", [])
    for chart_data in charts:
        _add_chart(slide, chart_data, theme)
    
    # Add each element from layout
    for elem in elements:
        elem_type = elem.get("type", "")
        x = elem.get("x", 0.5)
        y = elem.get("y", 0.5)
        width = elem.get("width", 5)
        height = elem.get("height", 1)
        style = elem.get("style", {})
        
        # Get content for this element type
        content = ""
        if elem_type in ["title", "heading"]:
            content = heading
        elif elem_type in ["subtitle", "subheading"]:
            content = subheading
        elif elem_type == "body":
            content = body_text
        elif elem_type == "bullets":
            content = bullets
        elif elem_type == "callout":
            content = callout
        
        # Add text element
        if content and elem_type in ["title", "heading", "subtitle", "subheading", "body"]:
            _add_text_element(slide, content, x, y, width, height, elem_type, style, theme)
        elif content and elem_type == "bullets":
            _add_bullet_element(slide, content, x, y, width, height, style, theme)
        elif content and elem_type == "callout":
            _add_callout_box(slide, content, theme, x, y, width, height)


def _add_text_element(
    slide,
    content: str,
    x: float,
    y: float,
    width: float,
    height: float,
    elem_type: str,
    style: Dict[str, Any],
    theme: Dict[str, Any]
) -> None:
    """Add a text element to the slide."""
    
    font_size = style.get("font_size", theme.get("heading_size", 32) if elem_type in ["title", "heading"] else theme.get("body_size", 18))
    alignment = style.get("alignment", "left")
    weight = style.get("weight", "bold" if elem_type in ["title", "heading"] else "normal")
    
    left = Inches(x)
    top = Inches(y)
    width_in = Inches(width)
    height_in = Inches(height)
    
    txBox = slide.shapes.add_textbox(left, top, width_in, height_in)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = content
    
    # Set font properties
    font = p.font
    font.size = Pt(font_size)
    font.name = theme.get("body_font", "Arial")
    font.bold = (weight == "bold")
    
    # Set alignment
    if alignment == "center":
        p.alignment = PP_ALIGN.CENTER
    elif alignment == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT


def _add_bullet_element(
    slide,
    bullets: List[str],
    x: float,
    y: float,
    width: float,
    height: float,
    style: Dict[str, Any],
    theme: Dict[str, Any]
) -> None:
    """Add bullet points to the slide."""
    
    font_size = style.get("font_size", theme.get("bullet_size", 16))
    
    left = Inches(x)
    top = Inches(y)
    width_in = Inches(width)
    height_in = Inches(height)
    
    txBox = slide.shapes.add_textbox(left, top, width_in, height_in)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = bullet
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.name = theme.get("body_font", "Arial")


def _add_default_content(slide, body_text: str, bullets: List[str], theme: Dict[str, Any]) -> None:
    """Add default content using placeholder."""
    
    # Get content placeholder
    content_placeholder = None
    for shape in slide.placeholders:
        if shape.placeholder_format.idx == 1:  # Body placeholder
            content_placeholder = shape
            break
    
    if content_placeholder:
        tf = content_placeholder.text_frame
        tf.clear()
        
        if body_text:
            p = tf.paragraphs[0]
            p.text = body_text
            p.level = 0
        
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
            p.level = 0
            p.font.size = Pt(theme.get("bullet_size", 16))


def _add_table(slide, table_data: Dict[str, Any], theme: Dict[str, Any]) -> None:
    """Add a table to the slide."""
    
    rows = table_data.get("rows", [])
    if not rows:
        return
    
    x = table_data.get("x", 1)
    y = table_data.get("y", 2)
    width = table_data.get("width", 10)
    height = table_data.get("height", 3)
    
    num_rows = len(rows)
    num_cols = len(rows[0]) if rows else 0
    
    if num_cols == 0:
        return
    
    left = Inches(x)
    top = Inches(y)
    width_in = Inches(width)
    height_in = Inches(height)
    
    # Add table
    table = slide.shapes.add_table(num_rows, num_cols, left, top, width_in, height_in).table
    
    # Fill in data
    for i, row in enumerate(rows):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            # Style header row
            if i == 0:
                cell = table.cell(i, j)
                cell.fill.solid()
                cell.fill.fore_color.rgb = theme.get("primary_color", RGBColor(102, 126, 234))


def _add_image(slide, image_path: str, theme: Dict[str, Any]) -> None:
    """Add an image to the slide."""
    
    if not os.path.exists(image_path):
        return
    
    # Default position and size
    x = 1
    y = 2
    width = 5
    height = 3
    
    try:
        slide.shapes.add_picture(
            image_path,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
    except Exception:
        pass


def _add_chart(slide, chart_data: Dict[str, Any], theme: Dict[str, Any]) -> None:
    """Add a chart to the slide (placeholder - requires additional libraries)."""
    
    # Note: Full chart support requires pptx.chart.data and additional imports
    # This is a placeholder for chart integration
    chart_type = chart_data.get("type", "bar")
    title = chart_data.get("title", "Chart")
    
    # For now, add a placeholder shape indicating chart location
    x = chart_data.get("x", 1)
    y = chart_data.get("y", 2)
    width = chart_data.get("width", 10)
    height = chart_data.get("height", 5)
    
    try:
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y),
            Inches(width), Inches(height)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.text_frame.text = f"[Chart: {title}]"
    except Exception:
        pass


def _apply_background(slide, background: Dict[str, Any], theme: Dict[str, Any]) -> None:
    """Apply background styling to the slide."""
    
    bg_type = background.get("type", "")
    
    if bg_type == "gradient":
        colors = background.get("colors", ["#667eea", "#764ba2"])
        # Apply gradient-like effect using shape
        try:
            background_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(7.5)
            )
            background_shape.fill.solid()
            # Convert hex to RGB
            if colors:
                primary_hex = colors[0].lstrip('#')
                r = int(primary_hex[0:2], 16)
                g = int(primary_hex[2:4], 16)
                b = int(primary_hex[4:6], 16)
                background_shape.fill.fore_color.rgb = RGBColor(r, g, b)
            background_shape.line.fill.background()
        except Exception:
            pass
    elif bg_type == "solid":
        color = background.get("color", "#212121")
        try:
            background_shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0), Inches(0),
                Inches(13.333), Inches(7.5)
            )
            background_shape.fill.solid()
            color_hex = color.lstrip('#')
            r = int(color_hex[0:2], 16)
            g = int(color_hex[2:4], 16)
            b = int(color_hex[4:6], 16)
            background_shape.fill.fore_color.rgb = RGBColor(r, g, b)
            background_shape.line.fill.background()
        except Exception:
            pass


def _add_callout_box(
    slide,
    text: str,
    theme: Dict[str, Any],
    x: float = 0.5,
    y: float = 5.5,
    width: float = 4,
    height: float = 1.5
) -> None:
    """Add a callout/quote box to a slide."""
    try:
        # Add a text box for the callout
        left = Inches(x)
        top = Inches(y)
        width_in = Inches(width)
        height_in = Inches(height)
        
        txBox = slide.shapes.add_textbox(left, top, width_in, height_in)
        tf = txBox.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = f'"{text}"'
        p.font.size = Pt(14)
        p.font.italic = True
        
        # Add a border/shape
        shape = slide.shapes.add_shape(
            1,  # msoShapeRectangle
            left, top, width_in, height_in
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.line.color.rgb = RGBColor(200, 200, 200)
        shape.text_frame = tf
    except Exception:
        pass


def _create_placeholder_pptx(output_path: str) -> str:
    """Create a placeholder PPTX file when python-pptx is not available."""
    # Create minimal valid PPTX structure
    # This is a minimal valid PPTX (actually just a placeholder message)
    with open(output_path, "w") as f:
        f.write("PPTX generation requires python-pptx package. Please install it.")
    return output_path


def get_presentation_stats(json_deck: Dict[str, Any]) -> Dict[str, Any]:
    """Get statistics about the presentation."""
    slides = json_deck.get("slides", [])
    
    total_bullets = sum(len(s.get("bullet_points", [])) for s in slides)
    has_calls = any(s.get("callout") for s in slides)
    has_tables = any(s.get("tables") for s in slides)
    has_images = any(s.get("images") for s in slides)
    has_charts = any(s.get("charts") for s in slides)
    
    return {
        "slide_count": len(slides),
        "total_bullets": total_bullets,
        "slides_with_calls": sum(1 for s in slides if s.get("callout")),
        "slides_with_tables": sum(1 for s in slides if s.get("tables")),
        "slides_with_images": sum(1 for s in slides if s.get("images")),
        "slides_with_charts": sum(1 for s in slides if s.get("charts")),
        "types": list(set(s.get("type", "content") for s in slides))
    }


def get_available_themes() -> List[str]:
    """Get list of available theme names."""
    return list(DEFAULT_THEMES.keys())


def get_theme(theme_name: str) -> Optional[Dict[str, Any]]:
    """Get theme configuration by name."""
    return DEFAULT_THEMES.get(theme_name)
